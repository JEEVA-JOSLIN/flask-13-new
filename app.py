# from flask import Flask, request, jsonify
# import json
# import base64
# import magic
# import os
# # Flask setup
# app = Flask(__name__)

# # Directory to save uploaded files temporarily
# UPLOAD_FOLDER = './uploads'
# os.makedirs(UPLOAD_FOLDER, exist_ok=True)
# app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# # MasterApp class encapsulates file processing logic
# class MasterApp:
#     def __init__(self):
#         self.magic = magic.Magic()

#     def process_file(self, file_path):
#         file_format = self.identify_file_format(file_path)
#         content = {}
#         if "text" in file_format.lower():
#             content = self.process_text(file_path)
#         elif "pdf" in file_format.lower():
#             content = self.process_pdf(file_path)
#         elif "word" in file_format.lower():
#             content = self.process_docx(file_path)
#         elif "powerpoint" in file_format.lower():
#             content = self.process_pptx(file_path)
#         else:
#             return {"error": "Unsupported file type."}
#         output_path = os.path.splitext(file_path)[0] + ".json"
#         self.save_to_json(content, output_path)
#         return content

#     def identify_file_format(self, file_path):
#         return self.magic.from_file(file_path)

#     def save_to_json(self, content, output_path):
#         with open(output_path, 'w', encoding='utf-8') as f:
#             json.dump(content, f, ensure_ascii=False, indent=4)

#     def process_text(self, file_path):
#         with open(file_path, 'r', encoding='utf-8') as file:
#             text_content = file.read()
#             return {
#                 "page_1": {
#                     "text": {"content": text_content},
#                     "recognized_text": {},
#                     "images": [],
#                     "tables": {}
#                 }
#             }

#     def process_pdf(self, file_path):
#         import fitz  # PyMuPDF
#         doc = fitz.open(file_path)
#         content = {}
#         for page_num in range(len(doc)):
#             page = doc[page_num]
#             text = page.get_text()
#             images = []
#             for img_index, img in enumerate(page.get_images(full=True)):
#                 xref = img[0]
#                 base_image = doc.extract_image(xref)
#                 img_bytes = base_image["image"]
#                 image_base64 = base64.b64encode(img_bytes).decode('utf-8')
#                 images.append({"base64": image_base64})
#             content[f"page_{page_num + 1}"] = {
#                 "text": {"content": text},
#                 "recognized_text": {},
#                 "images": images,
#                 "tables": {}
#             }
#         doc.close()
#         return content

#     def process_docx(self, file_path):
#         from docx import Document
#         doc = Document(file_path)
#         content = {}
#         page_counter = 1
#         for para in doc.paragraphs:
#             content[f"page_{page_counter}"] = {
#                 "text": {"content": para.text.strip()},
#                 "recognized_text": {},
#                 "images": [],
#                 "tables": {}
#             }
#             page_counter += 1
#         return content

#     def process_pptx(self, file_path):
#         from pptx import Presentation
#         presentation = Presentation(file_path)
#         content = {}
#         page_counter = 1
#         for slide in presentation.slides:
#             text_content = []
#             for shape in slide.shapes:
#                 if shape.has_text_frame:
#                     for para in shape.text_frame.paragraphs:
#                         text_content.append(para.text.strip())
#             content[f"page_{page_counter}"] = {
#                 "text": {"content": " ".join(text_content)},
#                 "recognized_text": {},
#                 "images": [],
#                 "tables": {}
#             }
#             page_counter += 1
#         return content

# # Initialize the MasterApp
# app.master_app = MasterApp()

# @app.route('/')
# def home():
#     return "Hello from Flask API on Azure!"

# @app.route('/process', methods=['POST'])
# def process_file():
#     if 'file' not in request.files:
#         return jsonify({"error": "No file provided"}), 400
#     file = request.files['file']
#     if file.filename == '':
#         return jsonify({"error": "No selected file"}), 400

#     # Save the uploaded file temporarily
#     file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
#     file.save(file_path)

#     # Process the file using the MasterApp class
#     content = app.master_app.process_file(file_path)
    
#     # Return the content in JSON format
#     return jsonify(content)
# if __name__ == "__main__":
#     app.run(debug=True)
import base64
import magic
import cv2
import sys
import torch
import fitz 
import pytesseract
from pytesseract import Output
from io import BytesIO
import numpy as np
from transformers import TrOCRProcessor, VisionEncoderDecoderModel
from docx import Document
from pdf2image import convert_from_bytes
from docx.document import Document as _Document
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.shapes.picture import Picture
from flask import Flask, request, jsonify
from collections import OrderedDict
from craft import CRAFT
from imagproc import resize_aspect_ratio, normalizeMeanVariance
from craft_utils import getDetBoxes, adjustResultCoordinates
processor_eng = TrOCRProcessor.from_pretrained('microsoft/trocr-large-stage1')
model_eng = VisionEncoderDecoderModel.from_pretrained('microsoft/trocr-large-stage1')
device = torch.device("cuda" if torch.cuda.is_available() else "cpu") 
model_eng.to(device)
app = Flask(__name__)
def copyStateDict(state_dict):
    """Copy state dictionary to remove module prefix if present."""
    if list(state_dict.keys())[0].startswith("module"):
        start_prefix = "module."
        new_state_dict = OrderedDict()
        for k, v in state_dict.items():
            name = k[len(start_prefix):]
            new_state_dict[name] = v
        return new_state_dict
    return state_dict
def detect_text(image, device,trained_model_path='craft_mlt_25k.pth'):
    net = CRAFT()
    net.load_state_dict(copyStateDict(torch.load(trained_model_path, map_location=device)))
    net.eval()
    net = net.to(device)
    img_resized, target_ratio, _ = resize_aspect_ratio(
        image, 
        square_size=1280, 
        interpolation=cv2.INTER_AREA
    )
    ratio_h = ratio_w = 1 / target_ratio
    x = normalizeMeanVariance(img_resized)
    x = torch.from_numpy(x).permute(2, 0, 1)    
    x = x.unsqueeze(0)  
    x = x.to(device)
    with torch.no_grad():
        y,_ = net(x)
    score_text = y[0,:,:,0].cpu().data.numpy()
    score_link = y[0,:,:,1].cpu().data.numpy()
    boxes, _ = getDetBoxes(score_text, score_link, text_threshold=0.7, link_threshold=0.4, low_text=0.4)
    boxes = adjustResultCoordinates(boxes, ratio_w, ratio_h)
    centers = [(np.mean(box[:, 0]), np.mean(box[:, 1])) for box in boxes]
    boxes_with_centers = list(zip(boxes, centers))
    boxes_with_centers.sort(key=lambda b: b[1][1])
    rows = []
    current_row = [boxes_with_centers[0]]
    for i in range(1, len(boxes_with_centers)):
        _, center = boxes_with_centers[i]
        _, prev_center = current_row[-1]
            
        if abs(center[1] - prev_center[1]) <= 10:
            current_row.append(boxes_with_centers[i])
        else:
            rows.append(current_row)
            current_row = [boxes_with_centers[i]]
    if current_row:
        rows.append(current_row)
    for row in rows:
        row.sort(key=lambda b: b[1][0])
    return rows
class MasterApp:
    def __init__(self):
        self.magic = magic.Magic()

    def detect_language(self, image):
        osd_data = pytesseract.image_to_osd(image, output_type=Output.DICT)
        script = osd_data.get('script', 'Unknown')
        orientation = osd_data.get('orientation', 'Unknown')
        return script, orientation

    def process_file(self, file_content):
        try:
            file_format = magic.Magic().from_buffer(file_content)
            content = {}
            if "image" in file_format.lower():
                nparr = np.frombuffer(file_content, np.uint8)
                image = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                content = self.process_image(image)
            elif "text" in file_format.lower():
                content = self.process_text(file_content)
            elif "pdf" in file_format.lower():
                content = self.process_pdf(file_content)
            elif "word" in file_format.lower():
                content = self.process_docx(file_content)
            elif "powerpoint" in file_format.lower():
                content = self.process_pptx(file_content)
            else:
                print("error: Unsupported file type.")
                return -1
            return content
        except Exception as e:
            print(f"Error processing file: {e}")
            return -1
    
    def perform_ocr(self,image):
        rows=detect_text(image,device)
        text=""
        ordered_words = []
        script_to_model = {
           # 'Latin': (model_eng, processor_eng),       
        }
        for row in rows:
            for box, _ in row:
                box = np.int32(box)
                x_coords = box[:, 0]
                y_coords = box[:, 1]
                xmin, xmax = int(np.min(x_coords)), int(np.max(x_coords))
                ymin, ymax = int(np.min(y_coords)), int(np.max(y_coords))
                word_image = image[ymin:ymax, xmin:xmax]
                ordered_words.append(word_image)
            
        for i,word in enumerate(ordered_words): 
            print(f"[INFO] Detecting the script/language for word {i}...")
            script, orientation = self.detect_language(image) 
            print(f"[INFO] Detected Script: {script}, Orientation: {orientation}°")
            if script not in script_to_model:
                print("[WARNING] Script not recognized or not supported. Defaulting to English OCR.")
               # model, processor = model_eng, processor_eng
            else:
              #  model, processor = script_to_model[script]

            print(f"[INFO] Using OCR model for script: {script}")
            print("[INFO] Performing OCR...")
            pixel_values = processor(images=word, return_tensors="pt").pixel_values
           # generated_ids = model.generate(pixel_values)
            text += (processor.batch_decode(generated_ids, skip_special_tokens=True)[0]+" ")
            #print("\n",text,"\n")
        print("[INFO] OCR Result:")
        print(text)
        return text
    
    def process_image(self, image):
        ocr_text=self.perform_ocr(image)
        _, buffer = cv2.imencode('.png', image)
        image_base64 = base64.b64encode(buffer).decode('utf-8')
        print("ocr",ocr_text)
        return {
            "page_1": {
                "text": "",
                "recognized_text": [ocr_text],
                "images": [image_base64],
                "tables": []
            }
        }

    def process_text(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            text_content = file.read()

            return {
                "page_1": {
                    "text": text_content,
                    "recognized_text": [],
                    "images": [],
                    "tables": []
                }
            }
    
    def check_pdf_metadata(self,doc):
        metadata = doc.metadata
        for key, value in metadata.items():
            print(f"{key}: {value}")
        
        if "PowerPoint" in metadata.get('producer', '') or "PowerPoint" in metadata.get('creator', ''):
            return True
        else:
            return False
      
    def process_pdf(self, file):
        import shutil
        poppler_path = shutil.which("pdftoppm")
        pdf_file = BytesIO(file)
        doc = fitz.open(stream=pdf_file, filetype="pdf")
        is_powerpoint = self.check_pdf_metadata(doc)
        content = {}
        
        if is_powerpoint:
            pages = convert_from_bytes(file,poppler_path)
            for i, page in enumerate(pages):
                buffered = BytesIO()
                page.save(buffered, format="PNG")
                image_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
                
                image_bytes = base64.b64decode(image_base64)
                
                nparr = np.frombuffer(image_bytes, np.uint8)
                img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                ocr_text=self.perform_ocr(img)
                

                content[f"page_{i + 1}"] = {
                    "text": ocr_text,
                    "recognized_text": [],
                    "images": [image_base64],
                    "tables": []
                }
                
        else:
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text().strip()
                image_list = page.get_images(full=True)
                images=[]
                ocr_text=[]
                for img in image_list:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                    
                    image_np = np.frombuffer(image_bytes, np.uint8)
    
                    img = cv2.imdecode(image_np, cv2.IMREAD_COLOR) 
                    images.append(image_base64)
                    ocr_text.append(self.perform_ocr(img))

                content[f"page_{page_num + 1}"] = {
                    "text": text,
                    "recognized_text": ocr_text,
                    "images": images,
                    "tables": []
                }  
        doc.close()
        return content

    def iter_block_items(self,parent):
        if isinstance(parent, _Document):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            raise ValueError("something's not right")

        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    def extract_table_content(self,table):
        table_data = []
        for row in table.rows:
            row_data = [cell.text.strip() for cell in row.cells]
            table_data.append(row_data)
        return table_data

    def process_docx(self, file):
        docx_file = BytesIO(file)
        doc = Document(docx_file)
        content = { 
            "text":"",
            "recognized_text": [],
            "images": [],
            "tables": []
        }
        for block in self.iter_block_items(doc):
            if isinstance(block, Paragraph):
                if(block.text.strip()!=""):
                    content["text"]+=(block.text.strip()+"\n")
            elif isinstance(block, Table):
                content["tables"].append(self.extract_table_content(block))
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                image_part = rel.target_part
                image=image_part.blob
                crop_info = getattr(image_part, 'crop', None)
                nparr = np.frombuffer(image, np.uint8)
                img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                if crop_info:
                    crop_left = crop_info.get('left', 0)
                    crop_top = crop_info.get('top', 0)
                    crop_right = crop_info.get('right', 0)
                    crop_bottom = crop_info.get('bottom', 0)
                    h, w = img.shape[:2]
                    crop_left_px = int(w * crop_left)
                    crop_top_px = int(h * crop_top)
                    crop_right_px = int(w * (1 - crop_right))
                    crop_bottom_px = int(h * (1 - crop_bottom))
                    img = img[
                        crop_top_px:crop_bottom_px, 
                        crop_left_px:crop_right_px
                    ]
                ocr_text=self.perform_ocr(img)
                image_base64 = base64.b64encode(image).decode('utf-8')

                content["images"].append(image_base64)
                content["recognized_text"].append(ocr_text)

        return {"page_1":content}

    def process_pptx(self, file):
        pptx_file = BytesIO(file)
        presentation = Presentation(pptx_file)
        content = {}
        tot_pg=0
        for slide in presentation.slides:
            tot_pg+=1
            text=""
            table_data=[]
            image_base64=[]
            ocr_text=[]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text_content = paragraph.text.strip()
                        if text_content:  
                            text+=text_content
                    text+="\n"
                if shape.has_table:
                    table = shape.table
                    table_data.append(self.extract_table_content(table))  
                if isinstance(shape, Picture):
                    image = shape.image.blob
                    image_np = np.frombuffer(image, np.uint8)  
                    img = cv2.imdecode(image_np, cv2.IMREAD_COLOR)  
                    crop_info = {
                        'left': shape.crop_left,
                        'top': shape.crop_top,
                        'right': shape.crop_right,
                        'bottom': shape.crop_bottom
                    }
                    if any(crop_info.values()):
                        h, w = img.shape[:2]
                        crop_left_px = int(w * crop_info['left'])
                        crop_top_px = int(h * crop_info['top'])
                        crop_right_px = int(w * (1 - crop_info['right']))
                        crop_bottom_px = int(h * (1 - crop_info['bottom']))
                        img = img[
                            crop_top_px:crop_bottom_px, 
                            crop_left_px:crop_right_px
                        ]
                    ocr_text.append(self.perform_ocr(img))              
                    image_base64.append(base64.b64encode(image).decode('utf-8'))

            content[f"page_{tot_pg}"] = {
                "text": text,
                "recognized_text": ocr_text,
                "images": image_base64,
                "tables": table_data
            }
        return content
@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    file_content = file.read()

    content = app.master_app.process_file(file_content)
    print(jsonify(content))
    if content == -1:
        return jsonify({"error": "Unsupported file type"}), 400
    return jsonify(content), 200
if __name__ == "__main__":
    app.master_app = MasterApp()
    app.run(debug=True)














